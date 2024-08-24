import os
import re
from typing import List, Tuple
from tqdm import tqdm
from pptx import Presentation
import logging

def clean_text_for_xml(text: str) -> str:
    return re.sub(r'[^\x09\x0A\x0D\x20-\x7F]', '', text)

async def extract_notes_from_presentation(presentation_path: str) -> List[Tuple[int, str]]:
    if not os.path.exists(presentation_path):
        logging.error(f"Presentation file not found: {presentation_path}")
        raise FileNotFoundError(f"Presentation file not found: {presentation_path}")

    if not os.access(presentation_path, os.R_OK):
        logging.error(f"No read permissions for file: {presentation_path}")
        raise PermissionError(f"No read permissions for file: {presentation_path}")

    try:
        prs = Presentation(presentation_path)
    except Exception as e:
        logging.error(f"Error loading presentation {presentation_path}: {str(e)}")
        raise RuntimeError(f"Failed to load presentation: {str(e)}") from e

    notes_list = []
    for i, slide in enumerate(tqdm(prs.slides, desc="Extracting notes")):
        try:
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_text = notes_slide.notes_text_frame.text if notes_slide.notes_text_frame else ""
                notes_text = clean_text_for_xml(notes_text)
                notes_list.append((i + 1, notes_text if notes_text else "No notes found."))
            else:
                notes_list.append((i + 1, "No notes slide."))
        except Exception as e:
            logging.warning(f"Error processing slide {i + 1}: {str(e)}")
            notes_list.append((i + 1, f"Error processing slide: {str(e)}"))

    if not notes_list:
        logging.warning(f"No notes extracted from presentation: {presentation_path}")

    return notes_list