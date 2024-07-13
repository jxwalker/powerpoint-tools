import argparse
import re
import traceback
import logging
import json
import os
from pptx import Presentation
from docx import Document
from fpdf import FPDF
from ibm_watson import NaturalLanguageUnderstandingV1
from ibm_watson.natural_language_understanding_v1 import Features, SummarizationOptions
from ibm_cloud_sdk_core.authenticators import IAMAuthenticator
from typing import List
from tqdm import tqdm
from getpass import getpass

logging.basicConfig(level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

def load_config(config_path: str) -> dict:
    """Load configuration from a JSON file."""
    with open(config_path, 'r') as file:
        config = json.load(file)
    return config

def extract_notes_from_presentation(presentation_path: str) -> List[str]:
    """Extract and clean notes from each slide in the PowerPoint presentation."""
    try:
        prs = Presentation(presentation_path)  # Load the PowerPoint presentation
    except Exception as e:
        logging.error(f"Error loading presentation: {e}")
        raise

    notes_list = []  # Initialize an empty list to store notes

    # Iterate through each slide in the presentation
    for i, slide in enumerate(prs.slides):
        if slide.has_notes_slide:  # Check if the slide has notes
            notes_slide = slide.notes_slide  # Get the notes slide
            notes_text_frame = notes_slide.notes_text_frame  # Get the text frame of the notes
            if notes_text_frame:
                notes_text = notes_text_frame.text  # Extract the notes text
                notes_text = clean_text_for_xml(notes_text)  # Clean the notes text
                notes_list.append(notes_text if notes_text else "No notes found.")
            else:
                notes_list.append("No notes found.")
        else:
            notes_list.append("No notes slide.")
    
    return notes_list  # Return the list of notes

def analyze_text(text: str, summarization_level: int, api_key: str, service_url: str) -> str:
    """Analyze and summarize the provided text using IBM Watson."""
    # Set up IBM Watson NLU credentials and service
    authenticator = IAMAuthenticator(api_key)
    natural_language_understanding = NaturalLanguageUnderstandingV1(
        version='2022-04-07',
        authenticator=authenticator
    )
    natural_language_understanding.set_service_url(service_url)

    try:
        # Analyze the text with the specified summarization level
        response = natural_language_understanding.analyze(
            text=text,
            features=Features(summarization=SummarizationOptions(limit=summarization_level))
        ).get_result()
        if 'summarization' in response:
            return response['summarization']['text']  # Return the summarized text
        else:
            return "Summary not available"
    except Exception as e:
        logging.error(f"Error during Watson API call: {e}")
        traceback.print_exc()
        return "Summary not available"

def clean_text_for_xml(text: str) -> str:
    """Clean text to ensure it is XML compatible."""
    # Remove invalid XML characters
    return re.sub(r'[^\x09\x0A\x0D\x20-\x7F]', '', text)

def write_summary_to_word(summaries: List[str], notes: List[str], output_path: str, extract_only: bool) -> None:
    """Write summaries and notes to a Word document."""
    doc = Document()  # Create a new Word document
    doc.add_heading('Presentation Summary', 0)  # Add a main heading

    # Iterate through summaries and notes, adding them to the document
    for i, (summary, note) in enumerate(zip(summaries, notes)):
        doc.add_heading(f'Slide {i + 1}', level=1)  # Add slide heading
        if not extract_only and len(note) >= 100:  # Only add summary if not extracting only and notes are 100 characters or more
            doc.add_heading('Summary', level=2)  # Add summary heading
            doc.add_paragraph(summary)  # Add the summary text
        doc.add_heading('Notes', level=2)  # Add notes heading
        doc.add_paragraph(note)  # Add the original notes
    
    doc.save(output_path)  # Save the document to the specified output path
    logging.info(f'Summary document saved to {output_path}')

def write_summary_to_text(summaries: List[str], notes: List[str], output_path: str, extract_only: bool) -> None:
    """Write summaries and notes to a plain text file."""
    with open(output_path, 'w') as file:
        file.write('Presentation Summary\n')
        for i, (summary, note) in enumerate(zip(summaries, notes)):
            file.write(f'**Slide {i + 1}**\n')  # Bold heading for slide number
            if not extract_only and len(note) >= 100:
                file.write('Summary\n')
                file.write(summary + '\n')
            file.write('Notes\n')
            file.write(note + '\n')
    logging.info(f'Summary document saved to {output_path}')

def write_summary_to_markdown(summaries: List[str], notes: List[str], output_path: str, extract_only: bool) -> None:
    """Write summaries and notes to a Markdown file."""
    with open(output_path, 'w') as file:
        file.write('# Presentation Summary\n')
        for i, (summary, note) in enumerate(zip(summaries, notes)):
            file.write(f'## **Slide {i + 1}**\n')  # Bold heading for slide number
            if not extract_only and len(note) >= 100:
                file.write('### Summary\n')
                file.write(summary + '\n')
            file.write('### Notes\n')
            file.write(note + '\n')
    logging.info(f'Summary document saved to {output_path}')

def write_summary_to_pdf(summaries: List[str], notes: List[str], output_path: str, extract_only: bool) -> None:
    """Write summaries and notes to a PDF file."""
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Arial", size=12)
    pdf.cell(200, 10, txt="Presentation Summary", ln=True, align='C')

    for i, (summary, note) in enumerate(zip(summaries, notes)):
        pdf.set_font("Arial", style='B', size=12)
        pdf.cell(200, 10, txt=f'Slide {i + 1}', ln=True, align='L')
        if not extract_only and len(note) >= 100:
            pdf.set_font("Arial", style='B', size=12)
            pdf.cell(200, 10, txt='Summary', ln=True, align='L')
            pdf.set_font("Arial", size=12)
            pdf.multi_cell(0, 10, txt=summary)
        pdf.set_font("Arial", style='B', size=12)
        pdf.cell(200, 10, txt='Notes', ln=True, align='L')
        pdf.set_font("Arial", size=12)
        pdf.multi_cell(0, 10, txt=note)
    
    pdf.output(output_path)
    logging.info(f'Summary document saved to {output_path}')

def validate_output_path(output_path: str):
    """Validate that the output path is writable."""
    dir_path = os.path.dirname(output_path)
    if not dir_path:
        dir_path = "."
    if not os.access(dir_path, os.W_OK):
        logging.error(f"Output path {dir_path} is not writable.")
        raise PermissionError(f"Output path {dir_path} is not writable.")

def main():
    # Set up argument parsing
    parser = argparse.ArgumentParser(description='Extract notes from a PowerPoint presentation and summarize them.')
    parser.add_argument('presentation_path', type=str, help='Path to the PowerPoint presentation file')
    parser.add_argument('output_path', type=str, help='Path to the output file')
    parser.add_argument('--summarization_level', type=int, default=7, help='Level of summarization (3-10, default is 7)')
    parser.add_argument('--config', type=str, default='config.json', help='Path to configuration file (default: config.json)')
    parser.add_argument('--output_format', type=str, choices=['docx', 'txt', 'md', 'pdf'], default='docx', help='Output format (default: docx)')
    parser.add_argument('--extract_only', action='store_true', help='Extract only the notes without summarizing')
    parser.add_argument('--log_to_file', action='store_true', help='Log output to a file')
    parser.add_argument('--verbose', action='store_true', help='Increase output verbosity')
    args = parser.parse_args()

    if args.summarization_level < 3 or args.summarization_level > 10:
        logging.error("Summarization level must be between 3 and 10.")
        return

    if args.verbose:
        logging.getLogger().setLevel(logging.DEBUG)

    if args.log_to_file:
        logging.basicConfig(filename='powerpoint_extractor.log', level=logging.DEBUG if args.verbose else logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')

    try:
        validate_output_path(args.output_path)
    except PermissionError as e:
        logging.error(f"Output path validation failed: {e}")
        return

    # Load configuration from file
    config = load_config(args.config)
    api_key = config.get('api_key') or getpass(prompt='Enter your IBM Watson API key: ')
    service_url = config.get('service_url') or input('Enter your IBM Watson service URL: ')

    try:
        notes = extract_notes_from_presentation(args.presentation_path)
        cleaned_notes = [note for note in notes if note.strip()]
        if not cleaned_notes:
            logging.error("No valid notes found in the presentation.")
            return
        summarized_notes = []
        if not args.extract_only:
            for i, note in enumerate(tqdm(cleaned_notes, desc="Processing slides")):
                if len(note) < 100:
                    logging.info(f"Notes are less than 100 characters, skipping summarization for slide {i + 1}.")
                    summarized_notes.append("")
                else:
                    summary = analyze_text(note, args.summarization_level, api_key, service_url)
                    summarized_notes.append(summary)
        else:
            summarized_notes = [""] * len(cleaned_notes)

        if args.output_format == 'docx':
            write_summary_to_word(summarized_notes, cleaned_notes, args.output_path, args.extract_only)
        elif args.output_format == 'txt':
            write_summary_to_text(summarized_notes, cleaned_notes, args.output_path, args.extract_only)
        elif args.output_format == 'md':
            write_summary_to_markdown(summarized_notes, cleaned_notes, args.output_path, args.extract_only)
        elif args.output_format == 'pdf':
            write_summary_to_pdf(summarized_notes, cleaned_notes, args.output_path, args.extract_only)
    except FileNotFoundError as e:
        logging.error(f"File not found: {e}")
        traceback.print_exc()
    except KeyError as e:
        logging.error(f"Key error: {e}. The expected response structure was not found.")
        traceback.print_exc()
    except Exception as e:
        logging.error(f"An unexpected error occurred: {e}")
        traceback.print_exc()

if __name__ == '__main__':
    main()
